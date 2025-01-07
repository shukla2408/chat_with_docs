import os
import streamlit as st
import faiss
import numpy as np
from io import BytesIO
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain.chains import RetrievalQA
from langchain.memory import ConversationSummaryBufferMemory
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.docstore.in_memory import InMemoryDocstore
from langchain.chat_models import ChatOpenAI
from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Set API key for OpenAI
openai_api_key = os.getenv("OPENAI_API_KEYS")
if not openai_api_key:
    st.error("OpenAI API Key is not set. Please add it to your environment variables.")
    st.stop()

# Get current directory
current_directory = os.getcwd()

# Paths for input folder
DEFAULT_INPUT_FOLDER = os.path.join(current_directory, 'input_files')
os.makedirs(DEFAULT_INPUT_FOLDER, exist_ok=True)

# Constants
MAX_MESSAGES = 5

# Google Drive Downloader class
class GoogleDriveDownloader:
    def __init__(self):
        self.credentials = self.get_credentials()
        self.service = build('drive', 'v3', credentials=self.credentials)

    def get_credentials(self):
        """Generate credentials from environment variables."""
        try:
            service_account_info = {
                "type": os.getenv("TYPE"),
                "project_id": os.getenv("PROJECT_ID"),
                "private_key_id": os.getenv("PRIVATE_KEY_ID"),
                "private_key": os.getenv("PRIVATE_KEY").replace("\\n", "\n"),
                "client_email": os.getenv("CLIENT_EMAIL"),
                "client_id": os.getenv("CLIENT_ID"),
                "auth_uri": os.getenv("AUTH_URI"),
                "token_uri": os.getenv("TOKEN_URI"),
                "auth_provider_x509_cert_url": os.getenv("AUTH_PROVIDER_X509_CERT_URL"),
                "client_x509_cert_url": os.getenv("CLIENT_X509_CERT_URL"),
            }
            return Credentials.from_service_account_info(service_account_info)
        except Exception as e:
            st.error(f"Error loading Google Drive credentials: {e}")
            st.stop()

    def download_files_from_folder(self, folder_id, local_path):
        try:
            query = f"'{folder_id}' in parents and trashed=false"
            results = self.service.files().list(q=query, fields="files(id, name)").execute()
            files = results.get('files', [])

            if not files:
                st.warning("No files found in the specified Google Drive folder.")
                return

            os.makedirs(local_path, exist_ok=True)

            for file in files:
                self._download_file(file['id'], file['name'], local_path)

        except Exception as e:
            st.error(f"Error downloading files: {e}")

    def _download_file(self, file_id, file_name, local_path):
        try:
            file_path = os.path.join(local_path, file_name)
            request = self.service.files().get_media(fileId=file_id)

            with open(file_path, 'wb') as f:
                downloader = MediaIoBaseDownload(f, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
        except Exception as e:
            st.error(f"Error downloading file {file_name}: {e}")

# Functions for processing input and answering questions
def process_input(input_folder, batch_size=50):
    """Processes files in the input folder and returns a vectorstore."""
    texts = []

    for filename in os.listdir(input_folder):
        file_path = os.path.join(input_folder, filename)
        try:
            if filename.endswith(".pdf"):
                pdf_reader = PdfReader(file_path)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        texts.append(page_text)
            elif filename.endswith(".docx"):
                doc = Document(file_path)
                for para in doc.paragraphs:
                    if para.text:
                        texts.append(para.text)
            elif filename.endswith(".pptx"):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text)
            elif filename.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as txt_file:
                    texts.append(txt_file.read())
        except Exception as e:
            st.warning(f"Error processing file {filename}: {e}")

    if not texts:
        raise ValueError("No text content extracted from the provided files.")

    # Use RecursiveCharacterTextSplitter for efficient splitting
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_texts = text_splitter.split_text(" ".join(texts))

    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)

    # Initialize FAISS vector store
    sample_embedding = np.array(embeddings.embed_query("sample text"))
    dimension = sample_embedding.shape[0]
    index = faiss.IndexFlatL2(dimension)

    # Create a vectorstore with batch processing
    vector_store = FAISS(
        embedding_function=embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={},
    )

    # Batch embed texts to avoid memory overload
    for i in range(0, len(split_texts), batch_size):
        batch = split_texts[i:i+batch_size]
        vector_store.add_texts(batch)

    return vector_store

def answer_question(vectorstore, query, memory):
    """Answers a question based on the provided vectorstore and conversation history."""
    llm = ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key, temperature=0.7)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 2})  # Limit retrieval to top 2 docs
    qa = RetrievalQA.from_chain_type(
        llm=llm, 
        retriever=retriever, 
        memory=memory
    )
    response = qa.run(query)
    return response

def truncate_messages(messages, max_messages=MAX_MESSAGES):
    """Truncates the message list to the last `max_messages` messages, with a summary if needed."""
    if len(messages) > max_messages:
        summary = f"Previous conversation summarized: {' '.join(msg['content'] for msg in messages[:-max_messages])}"
        truncated = [{"role": "assistant", "content": summary}] + messages[-max_messages:]
        return truncated
    return messages

def main():
    st.title("RAG Q&A App with GPT-4")

    # Sidebar for selecting input source
    with st.sidebar:
        st.header("Choose Input Source")
        input_source = st.selectbox("Select an option:", ["Select", "Local Folder", "Google Drive"])

        if input_source == "Local Folder":
            input_folder = st.text_input("Enter the local folder path:", value=DEFAULT_INPUT_FOLDER, key="local_folder_path")
            if st.button("Process Local Files"):
                try:
                    with st.spinner("Processing local files..."):
                        vectorstore = process_input(input_folder)
                    st.session_state["vectorstore"] = vectorstore
                    st.session_state["memory"] = ConversationSummaryBufferMemory(
                        llm=ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key),
                        memory_key="chat_history", 
                        max_token_limit=500
                    )
                    st.session_state["messages"] = [
                        {"role": "assistant", "content": "Hello! How can I assist you today?"}
                    ]
                    st.success("Local files processed successfully!")
                except Exception as e:
                    st.error(f"Error processing local files: {e}")

        elif input_source == "Google Drive":
            folder_id = st.text_input("Enter Google Drive Folder ID:")
            if st.button("Download and Process Files from Google Drive"):
                if folder_id:
                    downloader = GoogleDriveDownloader()
                    try:
                        with st.spinner("Downloading files from Google Drive..."):
                            downloader.download_files_from_folder(folder_id, DEFAULT_INPUT_FOLDER)
                        with st.spinner("Processing downloaded files..."):
                            vectorstore = process_input(DEFAULT_INPUT_FOLDER)
                        st.session_state["vectorstore"] = vectorstore
                        st.session_state["memory"] = ConversationSummaryBufferMemory(
                            llm=ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key),
                            memory_key="chat_history", 
                            max_token_limit=500
                        )
                        st.session_state["messages"] = [
                            {"role": "assistant", "content": "Hello! How can I assist you today?"}
                        ]
                        st.success("Files downloaded and processed successfully!")
                    except Exception as e:
                        st.error(f"Error processing files: {e}")
                else:
                    st.error("Please enter a valid Google Drive Folder ID.")

    # Main chat interface
    if "vectorstore" in st.session_state and "memory" in st.session_state:
        if "messages" not in st.session_state:
            st.session_state["messages"] = [
                {"role": "assistant", "content": "Hello! How can I assist you today?"}
            ]

        st.session_state["messages"] = truncate_messages(st.session_state["messages"], MAX_MESSAGES)

        for msg in st.session_state.messages:
            st.chat_message(msg["role"]).write(msg["content"])

        if prompt := st.chat_input():
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.chat_message("user").write(prompt)

            st.session_state.messages = truncate_messages(st.session_state["messages"], MAX_MESSAGES)

            try:
                with st.spinner("Generating response..."):
                    vectorstore = st.session_state["vectorstore"]
                    memory = st.session_state["memory"]
                    response = answer_question(vectorstore, prompt, memory)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                    st.chat_message("assistant").write(response)
            except Exception as e:
                st.error(f"Error generating response: {e}")

if __name__ == "__main__":
    main()
